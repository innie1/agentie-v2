from __future__ import annotations

import csv
import io
import re
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt

from agentie.core.artifact_naming import artifact_filename,creator_from_session
from agentie.core.file_service import UPLOADS, inspect_file, unique_path
from agentie.core.memory_store import latest_assistant_text

_REFERENCE_RE = re.compile(r"\b(?:this|that|it|the previous answer|previous answer|last answer|above|what you just wrote|what you wrote)\b", re.I)
_DOCX_RE = re.compile(r"\b(?:create|make|generate|export|save|turn|convert)\b.*\b(?:docx|word document|word file)\b|\b(?:docx|word document|word file)\b.*\b(?:create|make|generate|export|save|turn|convert)\b", re.I)
_XLSX_RE = re.compile(r"\b(?:create|make|generate|export|save|turn|convert)\b.*\b(?:xlsx|excel|spreadsheet)\b|\b(?:xlsx|excel|spreadsheet)\b.*\b(?:create|make|generate|export|save|turn|convert)\b", re.I)
_PPTX_RE = re.compile(r"\b(?:create|make|generate|export|save|turn|convert)\b.*\b(?:pptx|powerpoint|presentation|slide deck|slides)\b|\b(?:pptx|powerpoint|presentation|slide deck|slides)\b.*\b(?:create|make|generate|export|save|turn|convert)\b", re.I)


def _safe_filename(name: str | None, suffix: str, prefix: str,creator:str="Agentie") -> str:
    return artifact_filename(creator,name,suffix,prefix)


def _extract_filename(message: str, suffix: str) -> str | None:
    m = re.search(rf"\b(?:called|named|as)\s+[\"']?([^\"']+?\{suffix})\b", message, re.I)
    return m.group(1).strip() if m else None


def _explicit_content(message: str) -> str | None:
    m = re.search(r"\b(?:with|using|from)\s+(?:the\s+)?(?:text|content|data)\s*[:\-]?\s*(.+)$", message, re.I | re.S)
    if not m: return None
    value = m.group(1).strip()
    if value and not _REFERENCE_RE.fullmatch(value.strip(" .?!\"'")): return value.strip(" \"'")
    return None


def _resolve_content(session_id: str, message: str) -> str | None:
    content = _explicit_content(message)
    if content is None and (_REFERENCE_RE.search(message) or len(message.split()) <= 14):content = latest_assistant_text(session_id, max_chars=120000)
    return content


def _plain(line: str) -> str:
    line = re.sub(r"^#{1,6}\s*", "", line.strip());line = re.sub(r"^[-*]\s+", "", line);line = re.sub(r"^\d+[.)]\s+", "", line);return re.sub(r"\*\*|__|`", "", line).strip()


def create_docx(content: str, filename: str | None = None,creator:str="Agentie") -> dict[str, Any]:
    if not content.strip(): raise ValueError("Word document content is empty.")
    UPLOADS.mkdir(parents=True, exist_ok=True); path = unique_path(_safe_filename(filename, ".docx", "document",creator));doc = Document(); title_used = False
    for raw in content.replace("\r\n", "\n").split("\n"):
        line = raw.strip()
        if not line: continue
        if line.startswith("# "):doc.add_heading(_plain(line), level=0 if not title_used else 1); title_used = True
        elif line.startswith("## "): doc.add_heading(_plain(line), level=1)
        elif line.startswith("### "): doc.add_heading(_plain(line), level=2)
        elif re.match(r"^[-*]\s+", line): doc.add_paragraph(_plain(line), style="List Bullet")
        elif re.match(r"^\d+[.)]\s+", line): doc.add_paragraph(_plain(line), style="List Number")
        else: doc.add_paragraph(_plain(line))
    if not doc.paragraphs: doc.add_paragraph(content)
    doc.core_properties.author = creator; doc.core_properties.title = path.stem.replace("-", " ").title(); doc.save(path);return inspect_file(path)


def _markdown_table(content: str) -> list[list[str]] | None:
    lines=[x.strip() for x in content.splitlines() if x.strip()]
    for i in range(len(lines)-1):
        if "|" not in lines[i] or not re.match(r"^\|?\s*:?-{3,}", lines[i+1]): continue
        rows=[]
        for line in lines[i:]:
            if "|" not in line: break
            cells=[c.strip() for c in line.strip("|").split("|")]
            if all(re.match(r"^:?-{3,}:?$", c.replace(" ","")) for c in cells): continue
            rows.append(cells)
        if rows:return rows
    return None


def _tabular_rows(content: str) -> list[list[str]]:
    table=_markdown_table(content)
    if table:return table
    try:
        sample=content[:4096]; dialect=csv.Sniffer().sniff(sample, delimiters=",\t;");rows=list(csv.reader(io.StringIO(content),dialect))
        if len(rows)>1 and max((len(r) for r in rows),default=0)>1:return rows[:5000]
    except Exception: pass
    lines=[_plain(x) for x in content.splitlines() if _plain(x)];return [["Content"], *[[x] for x in lines[:10000]]]


def create_xlsx(content: str, filename: str | None = None,creator:str="Agentie") -> dict[str, Any]:
    if not content.strip(): raise ValueError("Spreadsheet content is empty.")
    UPLOADS.mkdir(parents=True,exist_ok=True); path=unique_path(_safe_filename(filename,".xlsx","spreadsheet",creator));rows=_tabular_rows(content); wb=Workbook(); ws=wb.active; ws.title=creator[:31] or "Agentie"
    for r_idx,row in enumerate(rows,1):
        for c_idx,value in enumerate(row,1):
            cell=ws.cell(r_idx,c_idx,value=value);cell.alignment=Alignment(vertical="top",wrap_text=True)
            if r_idx==1: cell.font=Font(bold=True)
    for col in range(1,max((len(r) for r in rows),default=1)+1):
        width=min(50,max(12,max((len(str(ws.cell(r,col).value or "")) for r in range(1,min(ws.max_row,200)+1)),default=12)+2));ws.column_dimensions[get_column_letter(col)].width=width
    ws.freeze_panes="A2" if ws.max_row>1 else None; wb.save(path);return inspect_file(path)


def _sections(content: str) -> list[tuple[str,list[str]]]:
    title="Agentie Presentation"; sections=[]; current_title=None; bullets=[]
    for raw in content.replace("\r\n","\n").split("\n"):
        line=raw.strip()
        if not line: continue
        if line.startswith("# ") and title=="Agentie Presentation": title=_plain(line); continue
        if line.startswith(("## ","### ")):
            if current_title or bullets: sections.append((current_title or "Overview",bullets[:8]))
            current_title=_plain(line); bullets=[]
        else:
            text=_plain(line)
            if text: bullets.append(text)
    if current_title or bullets: sections.append((current_title or "Overview",bullets[:8]))
    if not sections:
        words=[x.strip() for x in re.split(r"(?<=[.!?])\s+",content) if x.strip()];sections=[("Overview",words[:6])]
    return [(title,[]),*sections[:18]]


def create_pptx(content: str, filename: str | None = None,creator:str="Agentie") -> dict[str, Any]:
    if not content.strip(): raise ValueError("Presentation content is empty.")
    UPLOADS.mkdir(parents=True,exist_ok=True); path=unique_path(_safe_filename(filename,".pptx","presentation",creator));prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); sections=_sections(content);title,_=sections[0]; slide=prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text=title; slide.placeholders[1].text=f"Generated by {creator} · Agentie"
    for heading,bullets in sections[1:]:
        slide=prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text=heading[:120]; frame=slide.placeholders[1].text_frame; frame.clear()
        for idx,b in enumerate(bullets or ["No additional details provided."]):
            p=frame.paragraphs[0] if idx==0 else frame.add_paragraph(); p.text=b[:500]; p.level=0; p.font.size=Pt(20)
    prs.save(path); return inspect_file(path)


def try_office_request(session_id: str, message: str) -> dict[str, Any] | None:
    kind=None
    if _DOCX_RE.search(message): kind="docx"
    elif _XLSX_RE.search(message): kind="xlsx"
    elif _PPTX_RE.search(message): kind="pptx"
    if not kind:return None
    content=_resolve_content(session_id,message)
    if not content:return {"message":f"What should I put in the {kind.upper()} file? Paste the content or say “use the previous answer.”","card":None,"needs_content":True}
    creator=creator_from_session(session_id)
    if kind=="docx":card=create_docx(content,_extract_filename(message,".docx"),creator)
    elif kind=="xlsx":card=create_xlsx(content,_extract_filename(message,".xlsx"),creator)
    else:card=create_pptx(content,_extract_filename(message,".pptx"),creator)
    return {"message":f"Created {card['name']}.","card":card,"needs_content":False}
